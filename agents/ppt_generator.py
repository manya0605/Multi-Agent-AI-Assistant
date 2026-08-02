from openai import OpenAI
from pptx import Presentation
from dotenv import load_dotenv
import os

load_dotenv()

client = OpenAI(
    api_key=os.getenv("OPENROUTER_API_KEY"),
    base_url="https://openrouter.ai/api/v1"
)

def generate_ppt(topic):

    prompt = f"""
Create a PowerPoint presentation on {topic}.

Generate exactly 6 slides in this format:

Slide 1: Title
- {topic}

Slide 2: Introduction
- Introduction to {topic}
- Basic definition
- Importance

Slide 3: Key Concepts
- Main concept 1
- Main concept 2
- Main concept 3

Slide 4: Advantages
- Advantage 1
- Advantage 2
- Advantage 3

Slide 5: Applications
- Application 1
- Application 2
- Application 3

Slide 6: Conclusion
- Summary
- Future scope
- Final thoughts

Keep every slide short with bullet points.
"""

    response = client.chat.completions.create(
        model="openrouter/free",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ]
    )

    ppt_text = response.choices[0].message.content

    prs = Presentation()

    slides = ppt_text.split("Slide ")

    for slide in slides:

        if slide.strip() == "":
            continue

        layout = prs.slide_layouts[1]
        ppt_slide = prs.slides.add_slide(layout)

        lines = slide.strip().split("\n")

        ppt_slide.shapes.title.text = lines[0]

        body = ppt_slide.placeholders[1].text_frame
        body.clear()

        for line in lines[1:]:

            line = line.strip()

            if line:
                p = body.add_paragraph()
                p.text = line.replace("-", "").strip()

    filename = f"{topic.replace(' ', '_')}.pptx"

    prs.save(filename)

    return filename